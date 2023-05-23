from pptx import Presentation
from pptx.util import Inches

# 创建一个Presentation对象
presentation = Presentation()

# 添加一个带有标题和内容的幻灯片
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "欢迎来到我的演示文稿"
content.text = "这是我的第一个幻灯片"

# 添加一个带有标题和内容的幻灯片
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "第二个幻灯片"
content.text = "这是我的第二个幻灯片"

# 保存演示文稿
presentation.save("my_presentation.pptx")
